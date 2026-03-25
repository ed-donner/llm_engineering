import os
import fitz
from docx import Document
import pytesseract
from PIL import Image
from pptx import Presentation
import pandas as pd
import glob

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"


RAW_BASE = "knowledge_base_raw"
MD_BASE = "knowledge_base_markdown" 


def pdf_to_md(path):
    doc = fitz.open(path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text


def docx_to_md(path):
    doc = Document(path)
    return "\n\n".join(p.text for p in doc.paragraphs)


def image_to_md(path):
    img = Image.open(path)
    return pytesseract.image_to_string(img)

def pptx_to_md(path):
    prs = Presentation(path)
    md = []

    for i, slide in enumerate(prs.slides, start=1):
        md.append(f"# Slide {i}\n")

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    md.append(text + "\n")

    return "\n".join(md)


def excel_or_csv_to_md(path):
    ext = os.path.splitext(path)[1].lower()
    md = []

    if ext in [".xlsx", ".xls"]:
        sheets = pd.read_excel(path, sheet_name=None)
        for sheet_name, df in sheets.items():
            md.append(f"# Sheet: {sheet_name}\n")
            md.append(df.to_markdown(index=False))
            md.append("\n")
    elif ext == ".csv":
        df = pd.read_csv(path)
        md.append("# CSV\n")
        md.append(df.to_markdown(index=False))
        md.append("\n")

    return "\n".join(md)


def convert_and_store(file_path):

    ext = os.path.splitext(file_path)[1].lower()
    folder = os.path.basename(os.path.dirname(file_path))

    if ext == ".pdf":
        text = pdf_to_md(file_path)

    elif ext == ".docx":
        text = docx_to_md(file_path)

    elif ext in [".png", ".jpg", ".jpeg"]:
        text = image_to_md(file_path)

    elif ext == ".pptx":
        text = pptx_to_md(file_path)

    elif ext in [".xlsx", ".xls", ".csv"]:
        text = excel_or_csv_to_md(file_path)

    else:
        print(f"Skipping unsupported file: {file_path}")
        return

    filename = os.path.basename(file_path)
    name_without_ext = os.path.splitext(filename)[0]

    md_folder = os.path.join(MD_BASE, folder)
    os.makedirs(md_folder, exist_ok=True)

    md_path = os.path.join(md_folder, name_without_ext + ".md")

    with open(md_path, "w", encoding="utf-8") as f:
        f.write(text)

    print(f"Converted -> {md_path}")


def ingest_knowledge_base():

    files = glob.glob(f"{RAW_BASE}/**/*", recursive=True)

    for file_path in files:

        if os.path.isfile(file_path):
            convert_and_store(file_path)

    print("Ingestion complete")


def analyze_markdown_kb():

    knowledge_base_path = "knowledge_base_markdown/**/*.md"
    files = glob.glob(knowledge_base_path, recursive=True)

    print(files)
    print(f"Found {len(files)} files in the knowledge base")

    entire_knowledge_base = ""

    for file_path in files:
        with open(file_path, 'r', encoding='utf-8') as f:
            entire_knowledge_base += f.read()
            entire_knowledge_base += "\n\n"

    print(f"Total characters in knowledge base: {len(entire_knowledge_base):,}")




if __name__ == "__main__":

    ingest_knowledge_base()
    analyze_markdown_kb()