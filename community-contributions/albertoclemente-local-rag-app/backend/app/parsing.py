"""
Document text extraction and parsing utilities.
Handles PDF, DOCX, TXT, MD, EPUB with OCR hook stub for future expansion.
"""

import re
import json
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple
import mimetypes

# Document processing libraries
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import ebooklib
    from ebooklib import epub
    EPUB_AVAILABLE = True
except ImportError:
    EPUB_AVAILABLE = False

from app.models import DocumentType
from app.settings import get_settings
from app.diagnostics import get_logger, performance_context
from app.markdown_converter import get_markdown_converter, MarkdownConversionError

logger = get_logger(__name__)


class ParseError(Exception):
    """Custom exception for parsing errors"""
    pass


class DocumentParser:
    """Handles text extraction from various document formats"""
    
    def __init__(self):
        self.settings = get_settings()
        self.markdown_converter = get_markdown_converter()
    
    async def parse_document(self, file_path: Path, doc_type: DocumentType) -> Dict[str, Any]:
        """
        Parse a document and extract text content with metadata.
        Attempts to convert to Markdown first using Docling for better structure preservation.
        
        Args:
            file_path: Path to the document file
            doc_type: Type of document to parse
            
        Returns:
            Dictionary containing parsed content and metadata
        """
        with performance_context("parse_document", doc_type=doc_type.value):
            logger.info(f"Parsing document: {file_path.name} (type: {doc_type.value})")
            
            # Try Docling conversion first if available and format is supported
            if self.markdown_converter.is_available():
                supported_formats = self.markdown_converter.get_supported_formats()
                file_ext = file_path.suffix.lstrip('.').lower()
                
                if file_ext in supported_formats:
                    try:
                        logger.info(f"Attempting Docling conversion to Markdown for {file_path.name}")
                        md_result = await self.markdown_converter.convert_to_markdown(file_path)
                        
                        # Convert Docling result to our standard format
                        return {
                            'document_type': doc_type.value,
                            'total_pages': md_result.get('page_count', 0),
                            'total_chars': md_result['char_count'],
                            'total_words': md_result['word_count'],
                            'metadata': md_result['metadata'],
                            'full_text': md_result['markdown'],
                            'markdown': md_result['markdown'],  # Keep markdown version
                            'parsed_at': datetime.utcnow().isoformat(),
                            'structure': md_result['structure'],
                            'conversion_method': 'docling'
                        }
                    except MarkdownConversionError as e:
                        logger.warning(f"Docling conversion failed, falling back to legacy parser: {e}")
                    except Exception as e:
                        logger.warning(f"Unexpected error in Docling conversion, falling back: {e}")
            
            # Fallback to legacy parsers
            try:
                if doc_type == DocumentType.PDF:
                    return await self._parse_pdf(file_path)
                elif doc_type == DocumentType.DOCX:
                    return await self._parse_docx(file_path)
                elif doc_type == DocumentType.TXT:
                    return await self._parse_txt(file_path)
                elif doc_type == DocumentType.MD:
                    return await self._parse_markdown(file_path)
                elif doc_type == DocumentType.EPUB:
                    return await self._parse_epub(file_path)
                elif doc_type == DocumentType.HTML:
                    return await self._parse_html(file_path)
                elif doc_type == DocumentType.PPTX:
                    return await self._parse_pptx(file_path)
                else:
                    raise ParseError(f"Unsupported document type: {doc_type}")
                    
            except Exception as e:
                logger.error(f"Failed to parse document {file_path.name}: {e}")
                raise ParseError(f"Parsing failed: {str(e)}")
    
    async def _parse_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Parse PDF document"""
        if not PDF_AVAILABLE:
            raise ParseError("PyPDF2 not available for PDF parsing")
        
        pages = []
        metadata = {}
        
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                
                # Extract metadata
                if reader.metadata:
                    metadata = {
                        'title': reader.metadata.get('/Title', ''),
                        'author': reader.metadata.get('/Author', ''),
                        'subject': reader.metadata.get('/Subject', ''),
                        'creator': reader.metadata.get('/Creator', ''),
                        'producer': reader.metadata.get('/Producer', ''),
                        'creation_date': str(reader.metadata.get('/CreationDate', '')),
                    }
                
                # Extract text from each page
                for page_num, page in enumerate(reader.pages, 1):
                    try:
                        text = page.extract_text()
                        # Clean up the text
                        text = self._clean_text(text)
                        
                        if text.strip():  # Only add non-empty pages
                            pages.append({
                                'page_number': page_num,
                                'text': text,
                                'char_count': len(text)
                            })
                    except Exception as e:
                        logger.warning(f"Failed to extract text from page {page_num}: {e}")
                        # TODO: Add OCR fallback hook here
                        pages.append({
                            'page_number': page_num,
                            'text': '',
                            'char_count': 0,
                            'parse_error': str(e),
                            'needs_ocr': True
                        })
        
        except Exception as e:
            raise ParseError(f"PDF parsing error: {str(e)}")
        
        # Calculate total statistics
        total_text = '\n\n'.join(page['text'] for page in pages if page['text'])
        
        return {
            'document_type': 'pdf',
            'total_pages': len(pages),
            'total_chars': len(total_text),
            'total_words': len(total_text.split()),
            'metadata': metadata,
            'pages': pages,
            'full_text': total_text,
            'parsed_at': datetime.utcnow().isoformat(),
            'structure': self._analyze_structure(total_text)
        }
    
    async def _parse_docx(self, file_path: Path) -> Dict[str, Any]:
        """Parse DOCX document"""
        if not DOCX_AVAILABLE:
            raise ParseError("python-docx not available for DOCX parsing")
        
        try:
            doc = DocxDocument(file_path)
            
            # Extract paragraphs
            paragraphs = []
            for i, paragraph in enumerate(doc.paragraphs):
                text = paragraph.text.strip()
                if text:
                    paragraphs.append({
                        'paragraph_number': i + 1,
                        'text': text,
                        'style': paragraph.style.name if paragraph.style else None
                    })
            
            # Extract metadata from core properties
            metadata = {}
            if doc.core_properties:
                props = doc.core_properties
                metadata = {
                    'title': props.title or '',
                    'author': props.author or '',
                    'subject': props.subject or '',
                    'keywords': props.keywords or '',
                    'created': str(props.created) if props.created else '',
                    'modified': str(props.modified) if props.modified else '',
                    'last_modified_by': props.last_modified_by or ''
                }
            
            # Combine all text
            full_text = '\n\n'.join(p['text'] for p in paragraphs)
            full_text = self._clean_text(full_text)
            
            return {
                'document_type': 'docx',
                'total_paragraphs': len(paragraphs),
                'total_chars': len(full_text),
                'total_words': len(full_text.split()),
                'metadata': metadata,
                'paragraphs': paragraphs,
                'full_text': full_text,
                'parsed_at': datetime.utcnow().isoformat(),
                'structure': self._analyze_structure(full_text)
            }
            
        except Exception as e:
            raise ParseError(f"DOCX parsing error: {str(e)}")
    
    async def _parse_txt(self, file_path: Path) -> Dict[str, Any]:
        """Parse plain text document"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            text = None
            used_encoding = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        text = file.read()
                        used_encoding = encoding
                        break
                except UnicodeDecodeError:
                    continue
            
            if text is None:
                raise ParseError("Could not decode text file with any supported encoding")
            
            # Clean the text
            text = self._clean_text(text)
            
            # Split into lines for structure
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            
            return {
                'document_type': 'txt',
                'total_lines': len(lines),
                'total_chars': len(text),
                'total_words': len(text.split()),
                'encoding': used_encoding,
                'metadata': {},
                'lines': lines,
                'full_text': text,
                'parsed_at': datetime.utcnow().isoformat(),
                'structure': self._analyze_structure(text)
            }
            
        except Exception as e:
            raise ParseError(f"TXT parsing error: {str(e)}")
    
    async def _parse_markdown(self, file_path: Path) -> Dict[str, Any]:
        """Parse Markdown document"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
            
            # Extract headings for structure
            headings = []
            heading_pattern = r'^(#{1,6})\s+(.+)$'
            
            for line_num, line in enumerate(text.split('\n'), 1):
                match = re.match(heading_pattern, line.strip())
                if match:
                    level = len(match.group(1))
                    title = match.group(2).strip()
                    headings.append({
                        'level': level,
                        'title': title,
                        'line_number': line_num
                    })
            
            # Clean text (remove markdown syntax for full text)
            clean_text = self._clean_markdown(text)
            
            return {
                'document_type': 'md',
                'total_headings': len(headings),
                'total_chars': len(clean_text),
                'total_words': len(clean_text.split()),
                'metadata': {},
                'headings': headings,
                'raw_markdown': text,
                'full_text': clean_text,
                'parsed_at': datetime.utcnow().isoformat(),
                'structure': self._analyze_structure(clean_text, headings)
            }
            
        except Exception as e:
            raise ParseError(f"Markdown parsing error: {str(e)}")
    
    async def _parse_epub(self, file_path: Path) -> Dict[str, Any]:
        """Parse EPUB document"""
        if not EPUB_AVAILABLE:
            raise ParseError("ebooklib not available for EPUB parsing")
        
        try:
            book = epub.read_epub(file_path)
            
            # Extract metadata
            metadata = {
                'title': book.get_metadata('DC', 'title')[0][0] if book.get_metadata('DC', 'title') else '',
                'author': book.get_metadata('DC', 'creator')[0][0] if book.get_metadata('DC', 'creator') else '',
                'language': book.get_metadata('DC', 'language')[0][0] if book.get_metadata('DC', 'language') else '',
                'publisher': book.get_metadata('DC', 'publisher')[0][0] if book.get_metadata('DC', 'publisher') else '',
                'identifier': book.get_metadata('DC', 'identifier')[0][0] if book.get_metadata('DC', 'identifier') else '',
            }
            
            # Extract text from chapters
            chapters = []
            full_text_parts = []
            
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    try:
                        # Simple HTML to text conversion
                        content = item.get_content().decode('utf-8')
                        text = self._html_to_text(content)
                        text = self._clean_text(text)
                        
                        if text.strip():
                            chapters.append({
                                'id': item.get_id(),
                                'file_name': item.get_name(),
                                'text': text,
                                'char_count': len(text)
                            })
                            full_text_parts.append(text)
                            
                    except Exception as e:
                        logger.warning(f"Failed to extract text from EPUB item {item.get_id()}: {e}")
            
            full_text = '\n\n'.join(full_text_parts)
            
            return {
                'document_type': 'epub',
                'total_chapters': len(chapters),
                'total_chars': len(full_text),
                'total_words': len(full_text.split()),
                'metadata': metadata,
                'chapters': chapters,
                'full_text': full_text,
                'parsed_at': datetime.utcnow().isoformat(),
                'structure': self._analyze_structure(full_text)
            }
            
        except Exception as e:
            raise ParseError(f"EPUB parsing error: {str(e)}")
    
    async def _parse_html(self, file_path: Path) -> Dict[str, Any]:
        """Parse HTML document (fallback method - Docling is preferred)"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            
            # Simple HTML to text conversion
            text = self._html_to_text(html_content)
            text = self._clean_text(text)
            
            return {
                'document_type': 'html',
                'total_chars': len(text),
                'total_words': len(text.split()),
                'metadata': {},
                'full_text': text,
                'parsed_at': datetime.utcnow().isoformat(),
                'structure': self._analyze_structure(text),
                'conversion_method': 'legacy'
            }
        except Exception as e:
            raise ParseError(f"HTML parsing error: {str(e)}")
    
    async def _parse_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Parse PPTX document (fallback method - Docling is preferred)"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            slides_text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text_parts.append(text)
                
                if slide_text_parts:
                    slides_text.append({
                        'slide_number': slide_num,
                        'text': '\n'.join(slide_text_parts)
                    })
            
            full_text = '\n\n'.join(slide['text'] for slide in slides_text)
            full_text = self._clean_text(full_text)
            
            return {
                'document_type': 'pptx',
                'total_slides': len(slides_text),
                'total_chars': len(full_text),
                'total_words': len(full_text.split()),
                'metadata': {},
                'slides': slides_text,
                'full_text': full_text,
                'parsed_at': datetime.utcnow().isoformat(),
                'structure': self._analyze_structure(full_text),
                'conversion_method': 'legacy'
            }
        except Exception as e:
            raise ParseError(f"PPTX parsing error: {str(e)}")
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text content"""
        if not text:
            return ""
        
        # Normalize whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove excessive newlines
        text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
        
        # Strip leading/trailing whitespace
        text = text.strip()
        
        return text
    
    def _clean_markdown(self, text: str) -> str:
        """Remove markdown syntax to get plain text"""
        # Remove headers
        text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
        
        # Remove bold/italic
        text = re.sub(r'\*{1,2}([^*]+)\*{1,2}', r'\1', text)
        text = re.sub(r'_{1,2}([^_]+)_{1,2}', r'\1', text)
        
        # Remove links
        text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
        
        # Remove code blocks
        text = re.sub(r'```[^`]*```', '', text, flags=re.DOTALL)
        text = re.sub(r'`([^`]+)`', r'\1', text)
        
        # Remove lists
        text = re.sub(r'^\s*[-*+]\s+', '', text, flags=re.MULTILINE)
        text = re.sub(r'^\s*\d+\.\s+', '', text, flags=re.MULTILINE)
        
        return self._clean_text(text)
    
    def _html_to_text(self, html: str) -> str:
        """Simple HTML to text conversion"""
        # Remove script and style elements
        html = re.sub(r'<(script|style)[^>]*>.*?</\1>', '', html, flags=re.DOTALL | re.IGNORECASE)
        
        # Remove HTML tags
        text = re.sub(r'<[^>]+>', '', html)
        
        # Decode HTML entities
        import html as html_lib
        text = html_lib.unescape(text)
        
        return self._clean_text(text)
    
    def _analyze_structure(self, text: str, headings: Optional[List[Dict]] = None) -> Dict[str, Any]:
        """Analyze document structure for adaptive chunking hints"""
        if not text:
            return {'density': 'low', 'has_structure': False}
        
        # Calculate text density metrics
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        total_lines = len(lines)
        avg_line_length = sum(len(line) for line in lines) / max(total_lines, 1)
        
        # Detect structural elements
        has_headings = bool(headings) or bool(re.search(r'^[A-Z][^.!?]*$', text, re.MULTILINE))
        has_lists = bool(re.search(r'^\s*[-*•]\s+', text, re.MULTILINE))
        has_numbers = bool(re.search(r'^\s*\d+[\.)]\s+', text, re.MULTILINE))
        
        # Detect tables (simple heuristic)
        has_tables = bool(re.search(r'\|.*\|', text)) or text.count('\t') > 10
        
        # Determine density
        if avg_line_length > 80 and total_lines > 50:
            density = 'high'
        elif avg_line_length > 40 and total_lines > 20:
            density = 'medium'
        else:
            density = 'low'
        
        return {
            'density': density,
            'has_structure': has_headings or has_lists or has_numbers,
            'has_headings': has_headings,
            'has_lists': has_lists,
            'has_tables': has_tables,
            'avg_line_length': avg_line_length,
            'total_lines': total_lines,
            'heading_count': len(headings) if headings else 0
        }
    
    async def extract_text_for_ocr(self, file_path: Path) -> List[Dict[str, Any]]:
        """
        Extract images/pages that need OCR processing.
        This is a stub for future OCR integration.
        
        Returns:
            List of image data that needs OCR processing
        """
        # TODO: Implement OCR extraction
        # This would extract images from PDFs or scan documents
        # and return them in a format suitable for OCR processing
        
        logger.info(f"OCR extraction requested for: {file_path.name} (not implemented)")
        return []
    
    def get_supported_types(self) -> List[str]:
        """Get list of supported document types"""
        supported = ['txt', 'md', 'html', 'htm']
        
        if PDF_AVAILABLE:
            supported.append('pdf')
        if DOCX_AVAILABLE:
            supported.append('docx')
        if EPUB_AVAILABLE:
            supported.append('epub')
        
        # PPTX requires python-pptx
        try:
            import pptx
            supported.append('pptx')
        except ImportError:
            pass
        
        # Add Docling-supported formats if available
        if self.markdown_converter.is_available():
            docling_formats = self.markdown_converter.get_supported_formats()
            for fmt in docling_formats:
                if fmt not in supported:
                    supported.append(fmt)
        
        return supported
    
    def detect_file_type(self, file_path: Path) -> Optional[DocumentType]:
        """Detect document type from file"""
        # First try by extension
        ext = file_path.suffix.lower().lstrip('.')
        
        # Direct extension mapping
        ext_mapping = {
            'pdf': DocumentType.PDF,
            'docx': DocumentType.DOCX,
            'txt': DocumentType.TXT,
            'md': DocumentType.MD,
            'epub': DocumentType.EPUB,
            'html': DocumentType.HTML,
            'htm': DocumentType.HTML,
            'pptx': DocumentType.PPTX,
        }
        
        if ext in ext_mapping:
            return ext_mapping[ext]
        
        # Fall back to MIME type detection
        mime_type, _ = mimetypes.guess_type(str(file_path))
        if mime_type:
            mime_to_type = {
                'application/pdf': DocumentType.PDF,
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document': DocumentType.DOCX,
                'text/plain': DocumentType.TXT,
                'text/markdown': DocumentType.MD,
                'application/epub+zip': DocumentType.EPUB,
                'text/html': DocumentType.HTML,
                'application/vnd.openxmlformats-officedocument.presentationml.presentation': DocumentType.PPTX,
            }
            return mime_to_type.get(mime_type)
        
        return None


# Global parser instance
_parser_instance: Optional[DocumentParser] = None


def get_document_parser() -> DocumentParser:
    """Get the global document parser instance"""
    global _parser_instance
    if _parser_instance is None:
        _parser_instance = DocumentParser()
    return _parser_instance


async def get_document_parser_service() -> DocumentParser:
    """Get the global document parser service instance (async version)"""
    return get_document_parser()
